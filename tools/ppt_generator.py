from pptx import Presentation

def generate_ppt(path: str, title: str, content: str):
    """Generate a PowerPoint presentation (.pptx) with title and content."""
    try:
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = title
        slide.placeholders[1].text = content
        prs.save(path)
        return {"path": path, "status": "created"}
    except Exception as e:
        return {"error": str(e)}

# Tool export for dynamic loading
TOOL = {
    "name": "ppt_generator",
    "func": generate_ppt,
    "description": "Generate a PowerPoint presentation (.pptx) with title and content. Study platform use: create study presentations, lecture slides."
}
